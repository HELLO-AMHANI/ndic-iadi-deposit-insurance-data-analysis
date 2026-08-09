# BUILD SLIDES — slides deck (saved as .pptx then export to PDF)
# Author: Promise O. Amhanesi

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from PIL import Image as PILImage

FIGS = 'resultsfigs'
OUT  = 'docs/slides.pptx'

DARK_BLUE = RGBColor(0x1B, 0x3A, 0x6B)
MID_BLUE  = RGBColor(0x2E, 0x6B, 0xAD)
GREEN     = RGBColor(0x1D, 0x9E, 0x75)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF2, 0xF5, 0xF9)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=LIGHT_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def txb(slide, text, l, t, w, h,
        size=20, bold=False, color=DARK_BLUE,
        align=PP_ALIGN.LEFT, wrap=True):
    box  = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf   = box.text_frame
    tf.word_wrap = wrap
    p2   = tf.paragraphs[0]
    p2.alignment = align
    run  = p2.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return box

def add_img(slide, fname, l, t, w):
    path = os.path.join(FIGS, fname)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))
    else:
        txb(slide, f'[Chart: {fname}]', l, t, w, 1, size=11, color=MID_BLUE)

def footer(slide, text="Promise O. Amhanesi | NDIC + IADI Analysis | 2024–25"):
    rect(slide, 0, 7.1, 13.33, 0.4, fill=DARK_BLUE)
    txb(slide, text, 0.2, 7.15, 13, 0.35,
        size=9, color=WHITE, align=PP_ALIGN.LEFT)

def slide_header(slide, title, subtitle='', accent=MID_BLUE):
    rect(slide, 0, 0, 13.33, 1.1, fill=DARK_BLUE)
    txb(slide, title, 0.3, 0.1, 12, 0.55,
        size=26, bold=True, color=WHITE)
    if subtitle:
        txb(slide, subtitle, 0.3, 0.62, 12, 0.4,
            size=13, color=RGBColor(0xB5, 0xD4, 0xF4))

# SLIDE 1 — TITLE

s1 = add_slide()
bg(s1, DARK_BLUE)
txb(s1, "NDIC + IADI", 0.5, 1.2, 12.3, 1.0,
    size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s1, "Deposit Insurance Analysis", 0.5, 2.1, 12.3, 0.8,
    size=28, bold=False, color=RGBColor(0xB5, 0xD4, 0xF4), align=PP_ALIGN.CENTER)
txb(s1,
    "Benchmarking Nigeria's NDIC Against IADI Core Principles\n"
    "and Peer Member DICs (2010–2024)",
    0.5, 2.8, 12.3, 1.0,
    size=16, color=RGBColor(0xCC, 0xDD, 0xEE), align=PP_ALIGN.CENTER)
txb(s1, "Promise O. Amhanesi", 0.5, 4.5, 12.3, 0.5,
    size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(s1,
    "Banking & Finance · Philomath University Abuja | AMHANi Enterprise",
    0.5, 4.95, 12.3, 0.4,
    size=12, color=RGBColor(0xAA, 0xCC, 0xEE), align=PP_ALIGN.CENTER)
txb(s1,
    "github.com/HELLO-AMHANI/ndic-iadi-deposit-insurance-analysis",
    0.5, 5.5, 12.3, 0.4,
    size=11, color=RGBColor(0x1D, 0x9E, 0x75), align=PP_ALIGN.CENTER)

# SLIDE 2 — THE PROBLEM

s2 = add_slide()
bg(s2)
slide_header(s2, "1. The Problem", "Why benchmarking NDIC against IADI matters")
footer(s2)

bullets = [
    "Nigeria's NDIC protects depositors in 26 commercial banks but how robust is the fund?",
    "IADI publishes Core Principles for Effective Deposit Insurance: the global standard.",
    "No peer-reviewed quantitative benchmarking of NDIC against IADI exists for 2010–2024.",
    "This project fills that gap using Python, SQL, Excel, and machine learning.",
]
for i, b in enumerate(bullets):
    rect(s2, 0.5, 1.3 + i*1.3, 12.3, 1.1, fill=WHITE, line=MID_BLUE)
    txb(s2, f"→  {b}", 0.7, 1.4 + i*1.3, 11.8, 0.95, size=14)

# SLIDE 3 — IADI FRAMEWORK

s3 = add_slide()
bg(s3)
slide_header(s3, "2. IADI Core Principles Framework",
             "The 7 principles this project benchmarks against")
footer(s3)

cps = [
    ("CP1", "Public Policy Objectives"),
    ("CP2", "Mandate & Powers"),
    ("CP3", "Governance"),
    ("CP4", "Safety-Net Relationships"),
    ("CP5", "Membership & Coverage"),
    ("CP6", "Fund Management"),
    ("CP7", "Public Awareness"),
]
cols = [(0.4, 1.2), (4.5, 1.2), (8.6, 1.2),
        (0.4,  3.2), (4.5,  3.2), (8.6,  3.2),
        (4.5,  5.2)]
for (cp, label), (lx, ty) in zip(cps, cols):
    rect(s3, lx, ty, 3.8, 1.7, fill=DARK_BLUE)
    txb(s3, cp, lx+0.1, ty+0.1, 3.6, 0.6,
        size=22, bold=True, color=GREEN)
    txb(s3, label, lx+0.1, ty+0.65, 3.6, 0.9,
        size=12, color=WHITE)

# SLIDE 4 — DATA SOURCES

s4 = add_slide()
bg(s4)
slide_header(s4, "3. Data Sources & Methodology",
             "Cross-functional: Python · SQL · Excel · GitHub")
footer(s4)

srcs = [
    ("NDIC Annual Reports", "ndic.gov.ng/publications", "2010–2024 · 15 rows · manually extracted from PDF"),
    ("IADI Annual Survey",  "iadi.org",                "14 member countries · 15 years · 210 rows"),
    ("World Bank GFDD",     "databank.worldbank.org",   "GDP growth, inflation, FX rate, GDP per capita"),
    ("CBN Stat. Bulletin",  "cbn.gov.ng",               "Banking sector credit, monetary policy context"),
]
for i, (name, url, desc) in enumerate(srcs):
    rect(s4, 0.4, 1.3 + i*1.4, 12.5, 1.2, fill=WHITE, line=MID_BLUE)
    txb(s4, name, 0.6, 1.35 + i*1.4, 4, 0.5, size=13, bold=True, color=DARK_BLUE)
    txb(s4, url,  0.6, 1.75 + i*1.4, 4, 0.4, size=10, color=MID_BLUE)
    txb(s4, desc, 5.0, 1.5  + i*1.4, 7.7, 0.7, size=12, color=DARK_BLUE)

# SLIDE 5 — EDA HIGHLIGHTS

s5 = add_slide()
bg(s5)
slide_header(s5, "4. EDA Highlights",
             "Key patterns in NDIC fund performance 2010–2024")
footer(s5)
add_img(s5, '01_fund_vs_claims_timeseries.png', 0.4, 1.2, 7.5)
txb(s5, "Key Finding:", 8.2, 1.3, 4.9, 0.4, size=13, bold=True, color=GREEN)
findings = [
    "Fund grew 7.7× from ₦295bn (2010) to ₦2,279bn (2024)",
    "Claims flat at ~₦8.3bn per year from 2012–2023",
    "2024: claims spike to ₦63.2bn — Naira devaluation shock",
    "2016 oil crisis and 2020 COVID visible in trend",
    "Premium rate stable: 0.40–0.64% of total deposits",
]
for i, f in enumerate(findings):
    txb(s5, f"• {f}", 8.2, 1.8 + i*0.95, 4.9, 0.85, size=11.5)

# SLIDE 6 — BENCHMARKING SCORECARD

s6 = add_slide()
bg(s6)
slide_header(s6, "5. IADI Benchmarking Scorecard",
             "NDIC scored 1–5 against each IADI Core Principle")
footer(s6)
add_img(s6, '06_radar_iadi_benchmarking.png', 0.3, 1.1, 6.5)

scorecard = [
    ("CP1: Policy Objectives",    "4/5", "✅ Above avg"),
    ("CP2: Mandate & Powers",     "4/5", "✅ Above avg"),
    ("CP3: Governance",           "3/5", "⚠️ Below avg"),
    ("CP4: Safety-Net",           "4/5", "✅ Above avg"),
    ("CP5: Coverage",             "3/5", "🔴 Below avg"),
    ("CP6: Fund Management",      "4/5", "✅ Above avg"),
    ("CP7: Public Awareness",     "3/5", "⚠️ Below avg"),
]
for i, (cp, score, status) in enumerate(scorecard):
    rect(s6, 7.0, 1.2 + i*0.84, 5.9, 0.75,
         fill=WHITE if i%2==0 else LIGHT_BG, line=MID_BLUE)
    txb(s6, cp,     7.1, 1.28 + i*0.84, 3.2, 0.5, size=11)
    txb(s6, score,  10.3, 1.28 + i*0.84, 1.0, 0.5, size=11, bold=True, color=DARK_BLUE)
    txb(s6, status, 11.3, 1.28 + i*0.84, 1.5, 0.5, size=10, color=GREEN)

# SLIDE 7 — MODEL RESULTS

s7 = add_slide()
bg(s7)
slide_header(s7, "6. Predictive Modeling Results",
             "Fund stress classification train 2010–2019, test 2020–2024")
footer(s7)
add_img(s7, '08_roc_curve.png', 0.3, 1.1, 6.5)

txb(s7, "Model Performance", 7.1, 1.2, 5.8, 0.45, size=14, bold=True, color=DARK_BLUE)
model_rows = [
    ("Logistic Regression", "Baseline"),
    ("XGBoost", "Advanced"),
]
headers = ["Metric", "LR (Baseline)", "XGBoost"]
rect(s7, 7.1, 1.7, 5.8, 0.45, fill=DARK_BLUE)
for j, h_txt in enumerate(headers):
    txb(s7, h_txt, 7.1 + j*1.93, 1.75, 1.93, 0.38, size=11, bold=True, color=WHITE)

rows = [
    ("AUC-ROC", "—", "—"),
    ("Precision", "—", "—"),
    ("Recall", "—", "—"),
    ("LOO CV F1", "—", "—"),
]
for i, (m, lr_v, xgb_v) in enumerate(rows):
    bg_c = WHITE if i % 2 == 0 else LIGHT_BG
    rect(s7, 7.1, 2.15 + i*0.7, 5.8, 0.65, fill=bg_c, line=MID_BLUE)
    txb(s7, m,     7.1,    2.22 + i*0.7, 1.93, 0.5, size=11, color=DARK_BLUE)
    txb(s7, lr_v,  9.03,   2.22 + i*0.7, 1.93, 0.5, size=11, align=PP_ALIGN.CENTER)
    txb(s7, xgb_v, 11.0,   2.22 + i*0.7, 1.93, 0.5, size=11, align=PP_ALIGN.CENTER)

txb(s7, "⚠️ Small sample (n=14) — results are indicative, not definitive.",
    7.1, 4.9, 5.8, 0.5, size=10, color=RGBColor(0x99, 0x33, 0x00))
txb(s7, "Paste your actual metrics from results/model_metrics.csv into this slide.",
    7.1, 5.35, 5.8, 0.5, size=10, color=RGBColor(0x77, 0x77, 0x77))


# SLIDE 8 — SHAP INSIGHTS

s8 = add_slide()
bg(s8)
slide_header(s8, "7. SHAP Explainability",
             "Top predictors of fund stress machine learning transparency")
footer(s8)
add_img(s8, '11_shap_feature_importance.png', 0.3, 1.1, 7.0)

txb(s8, "Top 3 drivers of fund stress:", 7.6, 1.2, 5.4, 0.45,
    size=13, bold=True, color=DARK_BLUE)
insights = [
    ("1st", "Exchange Rate (lag 1yr)",
     "Naira depreciation in prior year\npredicts fund stress the next year"),
    ("2nd", "GDP Growth (lag 1yr)",
     "Economic recession = rising NPLs\nand claims pressure on NDIC"),
    ("3rd", "Inflation (lag 1yr)",
     "High inflation erodes depositor\nreal returns and bank profitability"),
]
for i, (rank, feat, desc) in enumerate(insights):
    rect(s8, 7.5, 1.75 + i*1.65, 5.5, 1.5, fill=WHITE, line=GREEN)
    txb(s8, rank, 7.6, 1.82 + i*1.65, 0.8, 0.5, size=18, bold=True, color=GREEN)
    txb(s8, feat, 8.4, 1.82 + i*1.65, 4.4, 0.5, size=12, bold=True, color=DARK_BLUE)
    txb(s8, desc, 7.6, 2.25 + i*1.65, 5.2, 0.85, size=11, color=DARK_BLUE)


# SLIDE 9 — LIMITATIONS

s9 = add_slide()
bg(s9)
slide_header(s9, "8. Limitations",
             "Transparent about what the data can and cannot tell us")
footer(s9)

lims = [
    ("Small sample", "15 annual observations results indicative, not statistically definitive"),
    ("Manual extraction", "NDIC data extracted from PDF reports minor variance possible"),
    ("2024 IADI divergence", "IADI-reported FAR (0.47) vs local (0.09) currency timing issue"),
    ("Loan ratio gap", "4 years show >5% LDR gap — NDIC uses net loans, data uses gross"),
    ("Peer coverage", "Only 14 of 82 IADI member countries available in public survey"),
]
for i, (title, desc) in enumerate(lims):
    rect(s9, 0.4, 1.3 + i*1.1, 12.4, 0.95, fill=WHITE, line=MID_BLUE)
    txb(s9, f"⚠️  {title}:", 0.6, 1.38 + i*1.1, 3.5, 0.5,
        size=12, bold=True, color=DARK_BLUE)
    txb(s9, desc, 4.2, 1.38 + i*1.1, 8.4, 0.6, size=12)


# SLIDE 10 — POLICY RECOMMENDATIONS + CONTACT

s10 = add_slide()
bg(s10)
slide_header(s10, "9–10. Policy Recommendations & Contact",
             "What NDIC and policymakers should do next")
footer(s10)

recs = [
    "Anchor FAR floor of 0.20 in statute to formalise fund adequacy targets",
    "Index coverage limit to GDP per capita annually prevent real erosion",
    "Publish annual FX stress test: fund adequacy at 30%, 50%, 70% devaluation",
    "Reform governance: fixed staggered board terms, supermajority removal rule",
    "Rural depositor awareness campaign especially microfinance bank clients",
]
for i, r in enumerate(recs):
    rect(s10, 0.4, 1.2 + i*0.95, 8.5, 0.82, fill=WHITE, line=GREEN)
    txb(s10, f"{i+1}.  {r}", 0.6, 1.28 + i*0.95, 8.1, 0.68, size=11.5)

rect(s10, 9.2, 1.2, 3.9, 5.5, fill=DARK_BLUE)
txb(s10, "Contact", 9.4, 1.4, 3.5, 0.5,
    size=15, bold=True, color=WHITE)
contact_items = [
    "Promise O. Amhanesi",
    "Banking & Finance",
    "Philomath University Abuja",
    "",
    "AMHANi Enterprise",
    "",
    "GitHub:",
    "HELLO-AMHANI/",
    "ndic-iadi-deposit-insurance-analysis",
]
for i, item in enumerate(contact_items):
    txb(s10, item, 9.4, 1.9 + i*0.48, 3.5, 0.45,
        size=10.5 if item else 6,
        color=GREEN if 'HELLO' in item or 'ndic' in item else WHITE)

prs.save(OUT)
print(f"✅  slides.pptx saved: {OUT}")
print("    Next step: Open in PowerPoint or LibreOffice → Export as PDF")
